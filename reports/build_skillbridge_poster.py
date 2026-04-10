from pathlib import Path
from typing import Optional

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent
OUTPUT_PNG = ROOT / "skillbridge_poster.png"
OUTPUT_PPTX = ROOT / "skillbridge_poster.pptx"

CANVAS_WIDTH = 4200
CANVAS_HEIGHT = 6000
MARGIN = 120

COLORS = {
    "white": "#FFFFFF",
    "navy": "#0E172A",
    "teal": "#0F766E",
    "sky": "#38BDF8",
    "coral": "#F97316",
    "gold": "#B59A57",
    "gold_light": "#E3D4A5",
    "stone": "#F5F1E8",
    "panel": "#FBFAF7",
    "text": "#1B1A18",
    "muted": "#706B63",
    "line": "#D8CCAB",
    "ink": "#121212",
}

ASSETS = {
    "ou_logo": ROOT / "images" / "Oakland_University_logo.png",
    "skillbridge_logo": ROOT / "images" / "newskillbridgelogo.png",
    "dashboard": ROOT / "images" / "FinalDashboard.png",
    "evidence": ROOT / "images" / "FinalEvidence.png",
    "job_match": ROOT / "images" / "FinalJobMatch.png",
    "erd": ROOT / "images" / "ERDDiagram.png",
    "core_flows": ROOT / "images" / "CoreFlowsDIagram.png",
}

FONTS = {
    "title": ("/System/Library/Fonts/Avenir Next.ttc", 0),
    "title_alt": ("/System/Library/Fonts/Avenir Next.ttc", 2),
    "body": ("/System/Library/Fonts/Avenir.ttc", 0),
    "body_bold": ("/System/Library/Fonts/Helvetica.ttc", 1),
}

HEADER_LINES = [
    "Group 14: Justin Elia, Jennifer Gonzalez, Spencer Roeren, Cordell Stonecipher",
    "CSI 4999: Senior Capstone Project - Winter 2026",
    "School of Engineering and Computer Science",
    "Oakland University",
    "Instructor: Dr. Hadeel M. Jawad",
]
TAGLINE = "Evidence-backed career intelligence platform"

DESCRIPTION = (
    "SkillBridge helps students and early-career job seekers turn resumes, projects, coursework, "
    "and other artifacts into a structured skill profile backed by real evidence. The platform "
    "supports resume ingestion, skill extraction and confirmation, job-fit analysis, skill-gap "
    "visibility, and tailored resume generation in one workflow."
)

AUDIENCE_LINE = (
    "Category: Career readiness and portfolio intelligence    "
    "Target users: students, interns, and early-career professionals"
)

STACK_CHIPS = [
    ("React + Vite", "sky"),
    ("FastAPI", "teal"),
    ("MongoDB", "gold"),
    ("Local AI", "coral"),
]

STAT_CARDS = [
    ("20", "Implemented workflows"),
    ("2", "Primary roles"),
    ("1", "Unified portfolio workspace"),
]

FEATURES = [
    "Resume ingestion, extraction review, and skill confirmation",
    "Evidence capture linked to projects and profile skills",
    "Job-match analysis with missing-skill and coverage reasoning",
    "Tailored resume generation with saved output history",
    "Admin moderation, role tagging, taxonomy, and analytics",
]


def hex_rgba(value: str, alpha: int = 255) -> tuple[int, int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4)) + (alpha,)


def load_font(kind: str, size: int) -> ImageFont.FreeTypeFont:
    path, index = FONTS[kind]
    return ImageFont.truetype(path, size=size, index=index)


def draw_shadow_card(
    image: Image.Image,
    box: tuple[int, int, int, int],
    *,
    radius: int = 36,
    fill: str = COLORS["white"],
    outline: Optional[str] = None,
    outline_width: int = 2,
    shadow_alpha: int = 38,
) -> None:
    overlay = Image.new("RGBA", image.size, (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(overlay)
    x0, y0, x1, y1 = box
    shadow_draw.rounded_rectangle(
        (x0 + 10, y0 + 18, x1 + 10, y1 + 18),
        radius=radius,
        fill=(14, 23, 42, shadow_alpha),
    )
    overlay = overlay.filter(ImageFilter.GaussianBlur(24))
    image.alpha_composite(overlay)

    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle(box, radius=radius, fill=hex_rgba(fill))
    if outline is not None:
        draw.rounded_rectangle(box, radius=radius, outline=hex_rgba(outline), width=outline_width)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = word if not current else f"{current} {word}"
        if draw.textlength(candidate, font=font) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    *,
    box: tuple[int, int, int, int],
    font: ImageFont.FreeTypeFont,
    fill: str,
    line_gap: int = 14,
) -> int:
    x0, y0, x1, y1 = box
    lines = wrap_text(draw, text, font, x1 - x0)
    line_height = font.getbbox("Ag")[3] - font.getbbox("Ag")[1]
    y = y0
    for line in lines:
        draw.text((x0, y), line, font=font, fill=fill)
        y += line_height + line_gap
        if y > y1:
            break
    return y


def paste_contain(
    canvas: Image.Image,
    source_path: Path,
    box: tuple[int, int, int, int],
    *,
    background: Optional[str] = None,
    rounded: Optional[int] = None,
    pad: int = 0,
) -> None:
    x0, y0, x1, y1 = box
    target_w = x1 - x0 - pad * 2
    target_h = y1 - y0 - pad * 2
    source = Image.open(source_path).convert("RGBA")
    contained = ImageOps.contain(source, (target_w, target_h), Image.Resampling.LANCZOS)

    layer = Image.new("RGBA", (target_w, target_h), (0, 0, 0, 0))
    if background:
        bg = Image.new("RGBA", (target_w, target_h), hex_rgba(background))
        layer.alpha_composite(bg)

    px = (target_w - contained.width) // 2
    py = (target_h - contained.height) // 2
    layer.alpha_composite(contained, (px, py))

    if rounded:
        mask = Image.new("L", (target_w, target_h), 0)
        mask_draw = ImageDraw.Draw(mask)
        mask_draw.rounded_rectangle((0, 0, target_w, target_h), radius=rounded, fill=255)
        canvas.paste(layer, (x0 + pad, y0 + pad), mask)
    else:
        canvas.alpha_composite(layer, (x0 + pad, y0 + pad))


def paste_cover(
    canvas: Image.Image,
    source_path: Path,
    box: tuple[int, int, int, int],
    *,
    rounded: int = 28,
) -> None:
    x0, y0, x1, y1 = box
    target_w = x1 - x0
    target_h = y1 - y0
    source = Image.open(source_path).convert("RGBA")
    covered = ImageOps.fit(source, (target_w, target_h), method=Image.Resampling.LANCZOS)
    mask = Image.new("L", (target_w, target_h), 0)
    mask_draw = ImageDraw.Draw(mask)
    mask_draw.rounded_rectangle((0, 0, target_w, target_h), radius=rounded, fill=255)
    canvas.paste(covered, (x0, y0), mask)


def draw_chip(
    image: Image.Image,
    *,
    x: int,
    y: int,
    text: str,
    bg: str,
    fg: str = COLORS["white"],
    height: int = 72,
    pad_x: int = 30,
) -> int:
    draw = ImageDraw.Draw(image)
    font = load_font("body_bold", 32)
    text_w = int(draw.textlength(text, font=font))
    width = text_w + pad_x * 2
    draw.rounded_rectangle((x, y, x + width, y + height), radius=height // 2, fill=hex_rgba(bg))
    text_y = y + (height - (font.getbbox("Ag")[3] - font.getbbox("Ag")[1])) // 2 - 4
    draw.text((x + pad_x, text_y), text, font=font, fill=fg)
    return width


def draw_section_heading(draw: ImageDraw.ImageDraw, text: str, x: int, y: int) -> None:
    font = load_font("title_alt", 50)
    text_w = int(draw.textlength(text, font=font))
    box = (x, y, x + text_w + 110, y + 84)
    draw.rounded_rectangle(box, radius=24, fill=hex_rgba(COLORS["navy"]))
    draw.rounded_rectangle((x + 18, y + 18, x + 34, y + 66), radius=8, fill=hex_rgba(COLORS["gold"]))
    draw.text((x + 54, y + 15), text, font=font, fill=COLORS["white"])


def draw_screenshot_card(
    canvas: Image.Image,
    *,
    box: tuple[int, int, int, int],
    title: str,
    image_path: Path,
) -> None:
    draw_shadow_card(canvas, box, radius=30, fill=COLORS["white"], outline=COLORS["line"])
    draw = ImageDraw.Draw(canvas)
    x0, y0, x1, y1 = box
    draw.rounded_rectangle((x0 + 18, y0 + 18, x1 - 18, y0 + 110), radius=24, fill=hex_rgba(COLORS["stone"]))
    draw.text((x0 + 36, y0 + 38), title, font=load_font("body_bold", 36), fill=COLORS["navy"])
    draw.rounded_rectangle((x0 + 36, y0 + 90, x0 + 250, y0 + 98), radius=4, fill=hex_rgba(COLORS["sky"]))
    image_box = (x0 + 24, y0 + 128, x1 - 24, y1 - 24)
    paste_contain(canvas, image_path, image_box, background=COLORS["white"], rounded=26, pad=12)


def draw_stat_card(
    canvas: Image.Image,
    *,
    box: tuple[int, int, int, int],
    number: str,
    label: str,
    accent: str,
) -> None:
    draw_shadow_card(canvas, box, radius=28, fill=COLORS["panel"], outline=COLORS["line"], shadow_alpha=28)
    draw = ImageDraw.Draw(canvas)
    x0, y0, x1, y1 = box
    draw.rounded_rectangle((x0 + 22, y0 + 22, x0 + 46, y1 - 22), radius=12, fill=hex_rgba(accent))
    draw.text((x0 + 86, y0 + 30), number, font=load_font("title_alt", 70), fill=COLORS["navy"])
    draw.text((x0 + 88, y0 + 122), label, font=load_font("body", 28), fill=COLORS["muted"])


def build_poster_png() -> None:
    for path in ASSETS.values():
        if not path.exists():
            raise FileNotFoundError(path)

    canvas = Image.new("RGBA", (CANVAS_WIDTH, CANVAS_HEIGHT), hex_rgba(COLORS["white"]))
    draw = ImageDraw.Draw(canvas)

    # Background atmosphere.
    atmosphere = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    atmosphere_draw = ImageDraw.Draw(atmosphere)
    atmosphere_draw.ellipse((-380, -140, 1240, 1080), fill=hex_rgba(COLORS["gold_light"], 76))
    atmosphere_draw.ellipse((3050, -120, 4580, 1180), fill=hex_rgba(COLORS["sky"], 34))
    atmosphere_draw.ellipse((3100, 5150, 4560, 6420), fill=hex_rgba(COLORS["teal"], 28))
    atmosphere = atmosphere.filter(ImageFilter.GaussianBlur(36))
    canvas.alpha_composite(atmosphere)

    # Header: match the template structure, then design the rest freely.
    top_band = (MARGIN, 110, CANVAS_WIDTH - MARGIN, 520)
    draw_shadow_card(canvas, top_band, radius=28, fill=COLORS["navy"], shadow_alpha=24)
    draw.rounded_rectangle((MARGIN, 110, CANVAS_WIDTH - MARGIN, 126), radius=8, fill=hex_rgba(COLORS["gold"]))

    paste_contain(canvas, ASSETS["ou_logo"], (170, 155, 980, 430), rounded=18, background=COLORS["white"], pad=18)
    paste_contain(canvas, ASSETS["skillbridge_logo"], (3270, 150, 3960, 450), background=COLORS["white"], rounded=18, pad=18)

    title_font = load_font("title_alt", 156)
    title_width = int(draw.textlength("SkillBridge", font=title_font))
    title_x = (CANVAS_WIDTH - title_width) // 2
    draw.text((title_x, 210), "SkillBridge", font=title_font, fill=COLORS["white"])
    tagline_font = load_font("body_bold", 46)
    tagline_width = int(draw.textlength(TAGLINE, font=tagline_font))
    tagline_x = (CANVAS_WIDTH - tagline_width) // 2
    draw.text((tagline_x, 380), TAGLINE, font=tagline_font, fill=hex_rgba(COLORS["gold_light"]))

    header_text_y = 610
    line_fonts = [
        load_font("body_bold", 34),
        load_font("body_bold", 33),
        load_font("body_bold", 33),
        load_font("body_bold", 33),
        load_font("body_bold", 33),
    ]
    for text, font in zip(HEADER_LINES, line_fonts):
        text_width = int(draw.textlength(text, font=font))
        draw.text(((CANVAS_WIDTH - text_width) // 2, header_text_y), text, font=font, fill=COLORS["text"])
        header_text_y += 53

    draw.rounded_rectangle((MARGIN, 930, CANVAS_WIDTH - MARGIN, 938), radius=4, fill=hex_rgba(COLORS["line"]))

    # Description.
    desc_box = (MARGIN, 990, CANVAS_WIDTH - MARGIN, 1860)
    draw_shadow_card(canvas, desc_box, radius=38, fill=COLORS["white"], outline=COLORS["line"])
    draw_section_heading(draw, "Project Description", 180, 1060)
    draw_wrapped_text(
        draw,
        DESCRIPTION,
        box=(180, 1175, 2500, 1640),
        font=load_font("body", 48),
        fill=COLORS["text"],
        line_gap=20,
    )

    chip_x = 180
    for label, color_key in STACK_CHIPS:
        width = draw_chip(canvas, x=chip_x, y=1660, text=label, bg=COLORS[color_key], height=78, pad_x=34)
        chip_x += width + 18

    stat_boxes = [
        (2740, 1150, 3340, 1410),
        (3370, 1150, 3970, 1410),
        (3050, 1440, 3650, 1700),
    ]
    stat_colors = [COLORS["teal"], COLORS["sky"], COLORS["gold"]]
    for box, (number, label), accent in zip(stat_boxes, STAT_CARDS, stat_colors):
        draw_stat_card(canvas, box=box, number=number, label=label, accent=accent)

    audience_box = (MARGIN, 1900, CANVAS_WIDTH - MARGIN, 2070)
    draw_shadow_card(canvas, audience_box, radius=30, fill=COLORS["stone"], outline=COLORS["line"], shadow_alpha=24)
    draw.text((170, 1942), AUDIENCE_LINE, font=load_font("body_bold", 38), fill=COLORS["navy"])

    # User interface section.
    ui_box = (MARGIN, 2130, CANVAS_WIDTH - MARGIN, 3870)
    draw_shadow_card(canvas, ui_box, radius=38, fill=COLORS["panel"], outline=COLORS["line"])
    draw_section_heading(draw, "User Interface", 180, 2200)
    draw.text(
        (180, 2320),
        "Key views from the final build show the guided flow from portfolio history to evidence capture and job-fit analysis.",
        font=load_font("body", 36),
        fill=COLORS["muted"],
    )

    screenshot_boxes = [
        (180, 2435, 1570, 3260),
        (1650, 2435, 2550, 3440),
        (2630, 2435, 4020, 3260),
    ]
    screenshot_specs = [
        ("Dashboard", ASSETS["dashboard"]),
        ("Evidence", ASSETS["evidence"]),
        ("Job Match", ASSETS["job_match"]),
    ]
    for box, (title, path) in zip(screenshot_boxes, screenshot_specs):
        draw_screenshot_card(canvas, box=box, title=title, image_path=path)

    # Database section.
    db_box = (MARGIN, 3930, 1970, 5860)
    draw_shadow_card(canvas, db_box, radius=38, fill=COLORS["white"], outline=COLORS["line"])
    draw_section_heading(draw, "Database", 180, 4010)
    draw.text(
        (180, 4135),
        "MongoDB collections normalize the entities behind profiles, evidence, skills, jobs, job-match runs, and tailored resumes.",
        font=load_font("body", 32),
        fill=COLORS["muted"],
    )
    paste_contain(canvas, ASSETS["erd"], (180, 4260, 1910, 5550), background=COLORS["white"], rounded=24, pad=14)
    draw.rounded_rectangle((180, 5575, 1910, 5790), radius=24, fill=hex_rgba(COLORS["stone"]))
    draw.text(
        (220, 5625),
        "Normalized schema supports identity, portfolio evidence, skill graphs, job ingestion, and saved tailored outputs.",
        font=load_font("body", 24),
        fill=COLORS["text"],
    )

    # Core flow section.
    flow_box = (2110, 3930, 4080, 4890)
    draw_shadow_card(canvas, flow_box, radius=38, fill=COLORS["white"], outline=COLORS["line"])
    draw_section_heading(draw, "System Flow", 2170, 4010)
    draw.text(
        (2170, 4135),
        "The platform separates user and admin workflows while keeping job analysis grounded in saved skills and proof.",
        font=load_font("body", 30),
        fill=COLORS["muted"],
    )
    paste_contain(canvas, ASSETS["core_flows"], (2170, 4230, 4020, 4820), background=COLORS["white"], rounded=22, pad=10)

    # Features section.
    feature_box = (2110, 5030, 4080, 5860)
    draw_shadow_card(canvas, feature_box, radius=38, fill=COLORS["navy"])
    draw.text((2170, 5110), "Implemented Scope", font=load_font("title_alt", 60), fill=COLORS["white"])
    draw.rounded_rectangle((2170, 5190, 2350, 5198), radius=4, fill=hex_rgba(COLORS["gold"]))

    bullet_font = load_font("body", 28)
    bullet_y = 5258
    for feature in FEATURES:
        draw.ellipse((2175, bullet_y + 10, 2199, bullet_y + 34), fill=hex_rgba(COLORS["gold"]))
        draw_wrapped_text(
            draw,
            feature,
            box=(2220, bullet_y, 4010, bullet_y + 96),
            font=bullet_font,
            fill="#F7F5F0",
            line_gap=8,
        )
        bullet_y += 88

    draw.rounded_rectangle((2170, 5720, 4010, 5808), radius=18, fill=hex_rgba("#FFFFFF", 22))
    draw.text(
        (2200, 5741),
        "Final build includes portfolio, analytics, moderation, onboarding, and resume-tailoring workflows.",
        font=load_font("body_bold", 26),
        fill=hex_rgba("#FFFFFF", 236),
    )

    OUTPUT_PNG.parent.mkdir(parents=True, exist_ok=True)
    canvas.convert("RGB").save(OUTPUT_PNG, quality=96)


def build_poster_pptx() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(28)
    presentation.slide_height = Inches(40)
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(OUTPUT_PNG), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    presentation.save(str(OUTPUT_PPTX))


def main() -> None:
    build_poster_png()
    build_poster_pptx()


if __name__ == "__main__":
    main()
