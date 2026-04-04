from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / ".pptx-export" / "slides"
POSTER_PATH = ROOT / "images" / "finalPresentationVideoPoster.png"
VIDEO_PATH = ROOT / "videos" / "final.presentation.demo.mp4"
OUTPUT_PATH = ROOT / "finalPresentation.pptx"


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slide_images = sorted(SLIDES_DIR.glob("slide-*.png"))
    if not slide_images:
        raise FileNotFoundError(f"No rendered slide PNGs found in {SLIDES_DIR}")

    for image_path in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    if not VIDEO_PATH.exists():
        raise FileNotFoundError(f"Missing demo video at {VIDEO_PATH}")
    if not POSTER_PATH.exists():
        raise FileNotFoundError(f"Missing poster frame at {POSTER_PATH}")

    video_slide = prs.slides.add_slide(blank_layout)
    video_slide.shapes.add_movie(
        str(VIDEO_PATH),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
        poster_frame_image=str(POSTER_PATH),
        mime_type="video/mp4",
    )

    prs.save(str(OUTPUT_PATH))


if __name__ == "__main__":
    main()
